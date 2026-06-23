import io
import logging
import asyncio
from typing import List
import re
from datetime import datetime

from groq import AsyncGroq, APIError as GroqAPIError
import google.generativeai as genai
from docx import Document
from pptx import Presentation

from config import GROQ_API_KEY, GEMINI_API_KEY, GROQ_MODEL, GEMINI_MODEL
from prompts import get_prompt

logger = logging.getLogger(__name__)

# Инициализация клиентов
groq_client = AsyncGroq(api_key=GROQ_API_KEY)

# Настройка Gemini
try:
    from google import genai as new_genai
    gemini_client = new_genai.Client(api_key=GEMINI_API_KEY)
    USE_NEW_GEMINI = True
    logger.info("✅ Используется новый google.genai клиент")
except ImportError:
    genai.configure(api_key=GEMINI_API_KEY)
    USE_NEW_GEMINI = False
    logger.info("✅ Используется старый google.generativeai клиент")

current_provider = "gemini"
MAX_CHUNK_SIZE = 2500
TIMEOUT_SECONDS = 30

gemini_failures = 0
last_gemini_fail_time = None
GROQ_FALLBACK_MINUTES = 5


# ============ ВСПОМОГАТЕЛЬНЫЕ ФУНКЦИИ ============

def split_text_into_chunks(text: str, max_size: int = MAX_CHUNK_SIZE) -> List[str]:
    """Разбивает текст на части по предложениям."""
    if len(text) <= max_size:
        return [text]
    
    chunks = []
    sentences = re.split(r'(?<=[.!?])\s+', text)
    current_chunk = ""
    
    for sentence in sentences:
        if len(current_chunk) + len(sentence) + 1 <= max_size:
            current_chunk += sentence + " "
        else:
            if current_chunk:
                chunks.append(current_chunk.strip())
            current_chunk = sentence + " "
    
    if current_chunk:
        chunks.append(current_chunk.strip())
    
    return chunks


def clean_markdown(text: str) -> str:
    """Удаляет Markdown-символы из текста."""
    text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
    text = re.sub(r'```[^\n]*\n.*?```', '', text, flags=re.DOTALL)
    text = re.sub(r'`([^`]+)`', r'\1', text)
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
    text = re.sub(r'\*([^*]+)\*', r'\1', text)
    text = re.sub(r'__([^_]+)__', r'\1', text)
    text = re.sub(r'_([^_]+)_', r'\1', text)
    text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'^[\s]*[-*+]\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'^[\s]*\d+\.\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'^[\s]*[-*_]{3,}\s*$', '', text, flags=re.MULTILINE)
    
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    text = '\n\n'.join(lines)
    
    return text

def remove_phone_mentions(text: str) -> str:
    """Удаляет упоминания телефонов, часов и посторонних устройств."""
    
    # 1. Удаляем по паттернам
    phone_patterns = [
        # Apple / iPhone / Watch
        r'iPhone\s*1S\s*',
        r'iPhone\s*',
        r'Apple\s*',
        r'iOS\s*',
        r'iCloud\s*',
        r'Mac\s*',
        r'iPad\s*',
        r'iPod\s*',
        r'Watch\s*',
        r'Apple\s*Watch\s*',
        r'умн[а-я]*\s*час[а-я]*\s*',
        
        # Общие устройства
        r'телефон[а-я]*\s*',
        r'смартфон[а-я]*\s*',
        r'гаджет[а-я]*\s*',
        r'мобильн[а-я]*\s*устройств[а-я]*\s*',
        r'устройств[а-я]*\s*',
        r'девайс[а-я]*\s*',
        r'планшет[а-я]*\s*',
        r'ноутбук[а-я]*\s*',
        r'компьютер[а-я]*\s*',
        
        # Действия с устройствами
        r'очистк[а-я]*\s*кэш[а-я]*\s*',
        r'очистк[а-я]*\s*памят[ии]\s*',
        r'очистк[а-я]*\s*данных\s*',
        r'удалени[ея]\s*ненужн[а-я]*\s*приложени[йя]\s*',
        r'закрыти[ея]\s*ненужн[а-я]*\s*приложени[йя]\s*',
        r'обновлени[ея]\s*операционн[а-я]*\s*систем[ы]\s*',
        r'обновлени[ея]\s*iOS\s*',
        r'настройк[а-я]*\s*телефон[а-я]*\s*',
        r'облачн[а-я]*\s*хранилищ[еа]\s*',
        r'установк[а-я]*\s*нов[а-я]*\s*чип[а-я]*\s*',
        r'чип[а-я]*\s*',
        
        # Сторонние приложения
        r'Clean\s*My\s*Photos\s*',
        r'Phone\s*Cleaner\s*',
        r'Cleaner\s*',
        r'Clean\s*My\s*',
        
        # Фразы про ускорение
        r'ускор[и-я]+ть\s*работ[уы]\s*',
        r'медленн[а-я]*\s*работ[а-я]*\s*',
        r'вернуть\s*исходн[а-я]*\s*скорост[ьи]\s*',
        r'способы\s*ускорения\s*работ[ы]\s*',
        r'ускорить\s*работ[уы]\s*',
    ]
    
    for pattern in phone_patterns:
        text = re.sub(pattern, '', text, flags=re.IGNORECASE)
    
    # 2. Удаляем целые предложения с упоминаниями устройств
    sentences = re.split(r'(?<=[.!?])\s+', text)
    filtered_sentences = []
    
    device_keywords = [
        'телефон', 'смартфон', 'iphone', 'apple', 'ios', 'icloud',
        'watch', 'гаджет', 'устройств', 'clean my', 'phone cleaner',
        'очистк', 'кэш', 'приложени', 'обновлени', 'настройк',
        'планшет', 'ноутбук', 'компьютер', 'девайс', 'чип',
        'память', 'данных', 'ускор', 'медленн', 'скорост',
        'батаре', 'заряд', 'изнаш', 'выключ', 'отключ'
    ]
    
    for sentence in sentences:
        sentence_lower = sentence.lower()
        # Проверяем, не содержит ли предложение ключевые слова об устройствах
        contains_device = False
        for keyword in device_keywords:
            if keyword in sentence_lower:
                contains_device = True
                break
        
        # Пропускаем только если в предложении нет упоминания 1С
        if contains_device and '1с' not in sentence_lower and '1с' not in sentence:
            continue  # Пропускаем это предложение
        
        filtered_sentences.append(sentence)
    
    text = '. '.join(filtered_sentences)
    
    # 3. Удаляем целые абзацы с упоминаниями устройств
    paragraphs = text.split('\n\n')
    cleaned_paragraphs = []
    
    for para in paragraphs:
        para_lower = para.lower()
        contains_device = False
        for keyword in device_keywords:
            if keyword in para_lower:
                contains_device = True
                break
        
        # Если абзац содержит упоминания устройств и нет 1С — пропускаем
        if contains_device and '1с' not in para_lower:
            continue
        
        cleaned_paragraphs.append(para)
    
    text = '\n\n'.join(cleaned_paragraphs)
    
    # 4. Убираем лишние пробелы и переводы строк
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # 5. Убираем двойные точки и запятые
    text = re.sub(r'\.\.+', '.', text)
    text = re.sub(r',,+', ',', text)
    
    # 6. Убираем пустые строки в начале и конце
    text = text.strip()
    
    return text


def remove_service_phrases(text: str) -> str:
    """Удаляет служебные фразы, которые не нужны в видео-сценарии."""
    
    service_phrases = [
        # Существующие фразы
        r'На этом данная лекция заканчивается[\.!]?\s*',
        r'Продолжим в следующем уроке[\.!]?\s*',
        r'Это конец лекции[\.!]?\s*',
        r'Спасибо за внимание[\.!]?\s*',
        r'До свидания[\.!]?\s*',
        r'На этом всё[\.!]?\s*',
        r'Это был последний слайд[\.!]?\s*',
        r'Время Улучшения Производительности 1S[\.!]?\s*',
        r'Рассмотрим четыре эффективных способа[\.!]?\s*',
        r'Ниже представлен обработанный текст[:.]?\s*',
        r'Редактированный текст[:.]?\s*',
        r'Предоставленный текст требует коррекции и улучшения стилистики[.:]?\s*',
        r'Обработанный текст[:.]?\s*',
        r'Результат обработки[:.]?\s*',
        r'Продолжим изучение темы в следующем уроке[\.!]?\s*',
        r'Продолжим в следующем уроке[\.!]?\s*',
        r'Вышеперечисленные методы и настройки могут существенно повысить производительность[\.!]?\s*',
        r'Регулярная оптимизация и настройка могут помочь достичь желаемого результата[\.!]?\s*',
        r'Давайте рассмотрим четыре основных способа[\.!]?\s*',
        r'Первый шаг к увеличению производительности[\.!]?\s*',
        r'Для повышения эффективности работы с 1С существуют различные методы[\.!]?\s*',
        r'Урок Заканчивается[\.!]?\s*',
        r'Урок заканчивается[\.!]?\s*',
        r'Окончательный результат будет после продолжения следующего урока[\.!]?\s*',
        r'Окончательный результат будет после[\.!]?\s*',
        r'Продолжение следующего урока[\.!]?\s*',
        r'В следующем уроке мы продолжим[\.!]?\s*',
        r'Это был последний слайд[\.!]?\s*',
        r'Вопрос: Как ускорить работу 1С\?[\.!]?\s*',
        r'Есть четыре способа, которые позволяют действительно ускорить работу с 1С[\.!]?\s*',
        r'В следующем разделе мы рассмотрим[\.!]?\s*',
        r'Сегодня мы разберемся в причинах[\.!]?\s*',
        r'Давайте разберем эти проблемы подробнее[\.!]?\s*',
    ]
    
    for phrase in service_phrases:
        text = re.sub(phrase, '', text, flags=re.IGNORECASE)
    
    # Убираем лишние переводы строк после удаления
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    return text


def remove_duplicate_ideas(text: str) -> str:
    """Удаляет дублирующиеся идеи и повторяющиеся формулировки."""
    
    # Разбиваем на абзацы
    paragraphs = text.split('\n\n')
    
    # Список ключевых идей, которые уже были озвучены
    seen_ideas = set()
    unique_paragraphs = []
    
    for para in paragraphs:
        para_clean = para.strip()
        if not para_clean:
            continue
        
        # Извлекаем ключевую идею абзаца (первые 50 символов)
        idea_key = para_clean[:50].lower()
        
        # Проверяем, не повторяет ли этот абзац уже сказанное
        is_duplicate = False
        for seen in seen_ideas:
            # Если абзац очень похож на уже виденный (>60% совпадения)
            if len(idea_key) > 20 and len(seen) > 20:
                words_idea = set(idea_key.split())
                words_seen = set(seen.split())
                if words_idea and words_seen:
                    similarity = len(words_idea & words_seen) / len(words_idea | words_seen)
                    if similarity > 0.6:
                        is_duplicate = True
                        break
        
        if not is_duplicate:
            seen_ideas.add(idea_key)
            unique_paragraphs.append(para_clean)
    
    return '\n\n'.join(unique_paragraphs)


def final_cleanup(text: str) -> str:
    """Финальная очистка текста."""
    
    # 0. Убираем упоминания телефонов
    text = remove_phone_mentions(text)
    
    # 1. Убираем служебные фразы
    text = remove_service_phrases(text)
    
    # 2. Приводим "DNS" к "1С" (если речь идет о 1С)
    text = re.sub(r'работы\s+DNS\s+с\s+одной\s+информационной\s+базой', 'работы 1С с одной информационной базой', text)
    text = re.sub(r'работы\s+DNS', 'работы 1С', text)
    text = re.sub(r'работу\s+DNS', 'работу 1С', text)
    text = re.sub(r'работа\s+DNS', 'работа 1С', text)
    text = re.sub(r'перевод\s+DNS', 'перевод 1С', text)
    text = re.sub(r'перехода\s+DNS', 'перехода 1С', text)
    text = re.sub(r'переход\s+DNS', 'переход 1С', text)
    text = re.sub(r'DNS\s+с\s+одной\s+информационной\s+базой', '1С с одной информационной базой', text)
    text = re.sub(r'DNS\s+с\s+одной\s+ИБ', '1С с одной ИБ', text)
    
    # 3. Приводим "1S" к "1С"
    text = re.sub(r'1S\.P5', '1С', text)
    text = re.sub(r'1S(?=[^а-яА-Яa-zA-Z])', '1С', text)
    text = re.sub(r'1S\.', '1С.', text)
    text = re.sub(r'1S\s', '1С ', text)
    
    # 4. Приводим "1SV" к "1С"
    text = re.sub(r'1SV', '1С', text)
    
    
    # 5. Убираем пробелы перед знаками препинания
    text = re.sub(r'\s+([.,!?;:])', r'\1', text)
    
    # 6. Убираем дублирующиеся вступления
    intro_patterns = [
        r'(Привет, друзья!?)\s*(Сегодня мы поговорим о том, когда нужно переводить)',
        r'(Сегодня мы поговорим о том, когда нужно переводить)\s*(Сегодня мы поговорим о том, когда нужно переводить)',
        r'(Давайте разберемся)\s*(Давайте разберемся)',
        r'(Всему виной блокировки)\s*(Всему виной блокировки)',
    ]
    
    for pattern in intro_patterns:
        text = re.sub(pattern, r'\1. \2', text, flags=re.IGNORECASE)
    
    # 7. Убираем дублирующиеся фразы в начале предложений
    lines = text.split('\n')
    seen_starts = set()
    cleaned_lines = []
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        start_key = line[:40].lower()
        
        if start_key in seen_starts:
            continue
        
        seen_starts.add(start_key)
        cleaned_lines.append(line)
    
    text = '\n'.join(cleaned_lines)
    
    # 8. Убираем дублирующиеся идеи
    text = remove_duplicate_ideas(text)
    
    # 9. Разбиваем на абзацы по смыслу
    sentences = re.split(r'(?<=[.!?])\s+', text)
    
    paragraphs = []
    current_paragraph = []
    current_length = 0
    
    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue
        
        sentence_len = len(sentence)
        
        if re.match(r'^(И|Но|А|Потому что|Так как|Однако|Поэтому|Также|Кроме того)\s', sentence, re.IGNORECASE):
            if current_paragraph:
                current_paragraph[-1] += ' ' + sentence
            else:
                current_paragraph.append(sentence)
            continue
        
        if current_length + sentence_len > 500 or sentence_len > 200:
            if current_paragraph:
                paragraphs.append(' '.join(current_paragraph))
                current_paragraph = []
                current_length = 0
        
        current_paragraph.append(sentence)
        current_length += sentence_len
    
    if current_paragraph:
        paragraphs.append(' '.join(current_paragraph))
    
    text = '\n\n'.join(paragraphs)
    
    # 10. Убираем повторяющиеся слова
    text = re.sub(r'\b(\w+)\s+\1\b', r'\1', text)
    
    # 11. Убираем лишние пробелы (но сохраняем разрывы абзацев "\n\n")
    text = re.sub(r'[ \t]+', ' ', text)
    text = re.sub(r' *\n *', '\n', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # 12. Убираем пустые строки в начале и конце
    text = text.strip()
    
    return text


def merge_and_clean_text(parts: List[str]) -> str:
    """Объединяет части и удаляет повторы."""
    if not parts:
        return ""
    
    # Объединяем все части
    full_text = '\n\n'.join(parts)
    
    # Удаляем маркдаун и скобки
    full_text = clean_markdown(full_text)
    full_text = re.sub(r'\[[^\]]*\]', '', full_text)
    
    # Удаляем служебные заголовки "=== Часть X ==="
    full_text = re.sub(r'=== Часть \d+ ===\n*', '', full_text)
    
    # Финальная очистка
    full_text = final_cleanup(full_text)
    
    return full_text


# ============ ОСНОВНЫЕ ФУНКЦИИ ============

async def process_with_gemini(text: str, mode: str) -> str:
    """Обработка текста через Gemini с таймаутом."""
    global gemini_failures, last_gemini_fail_time
    
    try:
        prompt = get_prompt(mode)
        
        async def make_request():
            if USE_NEW_GEMINI:
                response = gemini_client.models.generate_content(
                    model=GEMINI_MODEL,
                    contents=f"{prompt}\n\n{text}"
                )
                return response.text or "⚠️ Не удалось получить ответ от Gemini."
            else:
                model = genai.GenerativeModel(
                    model_name=GEMINI_MODEL,
                    system_instruction=prompt
                )
                response = await model.generate_content_async(text)
                return response.text or "⚠️ Не удалось получить ответ от Gemini."
        
        result = await asyncio.wait_for(make_request(), timeout=TIMEOUT_SECONDS)
        
        if gemini_failures > 0:
            gemini_failures = 0
            logger.info("✅ Gemini восстановлен")
        
        return result
        
    except asyncio.TimeoutError:
        logger.error(f"Gemini timeout after {TIMEOUT_SECONDS}s")
        gemini_failures += 1
        last_gemini_fail_time = datetime.now()
        raise Exception(f"Gemini timeout (>{TIMEOUT_SECONDS}s)")
        
    except Exception as e:
        logger.error(f"Gemini error: {e}")
        gemini_failures += 1
        last_gemini_fail_time = datetime.now()
        raise


async def process_with_groq(text: str, mode: str) -> str:
    """Обработка текста через Groq."""
    try:
        prompt = get_prompt(mode)
        
        response = await groq_client.chat.completions.create(
            model=GROQ_MODEL,
            max_tokens=4096,
            temperature=0.7,
            messages=[
                {"role": "system", "content": prompt},
                {"role": "user", "content": f"Обработай следующий текст:\n\n{text}"}
            ]
        )
        
        return response.choices[0].message.content or "⚠️ Не удалось получить ответ от Groq."
        
    except GroqAPIError as e:
        logger.error(f"Groq API error: {e}")
        raise
    except Exception as e:
        logger.error(f"Groq error: {e}")
        raise


def should_use_groq() -> bool:
    """Определяет, нужно ли использовать Groq вместо Gemini."""
    global gemini_failures, last_gemini_fail_time
    
    if gemini_failures >= 2:
        if last_gemini_fail_time:
            elapsed = (datetime.now() - last_gemini_fail_time).total_seconds() / 60
            if elapsed < GROQ_FALLBACK_MINUTES:
                logger.info(f"Используем Groq (Gemini упал {gemini_failures} раз, прошло {elapsed:.1f} мин)")
                return True
    
    return False


async def process_chunk(chunk: str, mode: str, provider: str) -> str:
    """Обрабатывает один чанк текста."""
    if provider == "gemini":
        return await process_with_gemini(chunk, mode)
    else:
        return await process_with_groq(chunk, mode)


async def process_text_with_ai(text: str, mode: str) -> str:
    """Обработка текста с автоматическим переключением и разбивкой на части."""
    global current_provider, gemini_failures
    
    if len(text) < 50:
        return "📝 Текст слишком короткий для обработки. Добавьте больше материала."
    
    chunks = split_text_into_chunks(text)
    total_chunks = len(chunks)
    
    if total_chunks > 1:
        logger.info(f"Текст разбит на {total_chunks} частей")
    
    use_groq = should_use_groq()
    if use_groq:
        current_provider = "groq"
    else:
        current_provider = "gemini"
    
    results = []
    force_groq = False
    
    for i, chunk in enumerate(chunks, 1):
        logger.info(f"Обработка части {i}/{total_chunks} (размер: {len(chunk)} символов)")
        
        provider_to_use = "groq" if force_groq else current_provider
        
        try:
            result = await process_chunk(chunk, mode, provider_to_use)
            clean_result = clean_markdown(result)
            clean_result = re.sub(r'\[[^\]]*\]', '', clean_result)
            clean_result = clean_result.strip()
            results.append(clean_result)
            
            if provider_to_use == "gemini":
                gemini_failures = 0
                
        except Exception as e:
            logger.error(f"Ошибка при обработке части {i} через {provider_to_use}: {e}")
            
            if provider_to_use == "gemini" and not force_groq:
                logger.info(f"Gemini упал, пробуем Groq для части {i}")
                try:
                    result = await process_chunk(chunk, mode, "groq")
                    clean_result = clean_markdown(result)
                    clean_result = re.sub(r'\[[^\]]*\]', '', clean_result)
                    clean_result = clean_result.strip()
                    results.append(clean_result)
                    force_groq = True
                    current_provider = "groq"
                    logger.info("Переключено на Groq для следующих частей")
                except Exception as e2:
                    logger.error(f"Groq тоже упал для части {i}: {e2}")
                    results.append(f"⚠️ Часть {i} не удалось обработать.")
            else:
                results.append(f"⚠️ Часть {i} не удалось обработать.")
        
        if i < total_chunks:
            await asyncio.sleep(1)
    
    # Объединяем и чистим
    full_result = merge_and_clean_text(results)
    
    return full_result


def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    """Извлекает текст из различных форматов файлов."""
    ext = filename.split('.')[-1].lower()
    
    try:
        if ext == 'txt':
            if isinstance(file_bytes, bytes):
                return file_bytes.decode('utf-8', errors='ignore')
            elif hasattr(file_bytes, 'getvalue'):
                return file_bytes.getvalue().decode('utf-8', errors='ignore')
            else:
                return str(file_bytes)
        
        elif ext == 'docx':
            if isinstance(file_bytes, bytes):
                doc = Document(io.BytesIO(file_bytes))
            elif hasattr(file_bytes, 'getvalue'):
                doc = Document(io.BytesIO(file_bytes.getvalue()))
            else:
                doc = Document(file_bytes)
            
            return '\n'.join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
        
        elif ext == 'pptx':
            if isinstance(file_bytes, bytes):
                prs = Presentation(io.BytesIO(file_bytes))
            elif hasattr(file_bytes, 'getvalue'):
                prs = Presentation(io.BytesIO(file_bytes.getvalue()))
            else:
                prs = Presentation(file_bytes)
            
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(shape.text)
            return '\n\n'.join(text)
        
        else:
            raise ValueError(f"Unsupported file format: {ext}")
            
    except Exception as e:
        logger.error(f"Error extracting text from {filename}: {e}")
        raise


async def process_file(file_bytes: bytes, filename: str, mode: str) -> str:
    """Обрабатывает загруженный файл."""
    try:
        text = extract_text_from_file(file_bytes, filename)
        
        if not text.strip():
            return "⚠️ Не удалось извлечь текст из файла. Файл пуст или содержит только изображения."
        
        result = await process_text_with_ai(text, mode)
        return result
        
    except Exception as e:
        logger.error(f"Error processing file: {e}")
        raise